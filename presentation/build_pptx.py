"""Build the figure-forward talk as an editable .pptx (16:9) for Google Slides.

Design rule: each slide = short title + large figure + AT MOST one short line.
No dense paragraphs, no duplicated captions. Restrained palette (blue = ours).
Figures are the real rendered PNGs in presentation/figures/ and assets/.
"""
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
ASSET = os.path.join(HERE, "assets")

# --- palette (matches the deck's restrained Okabe-Ito scheme) ----------------
INK = RGBColor(0x22, 0x22, 0x22)      # near-black title/body
MUTED = RGBColor(0x60, 0x60, 0x60)    # grey one-liners / slide numbers
BLUE = RGBColor(0x00, 0x72, 0xB2)     # "ours" accent (Okabe-Ito blue)

TITLE_FONT = "Segoe UI Semibold"
BODY_FONT = "Segoe UI"

EMU_PER_IN = 914400.0

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def fit(img_path, bl, bt, bw, bh):
    """Return (left, top, w, h) in inches fitting img into box, centered."""
    iw, ih = Image.open(img_path).size
    ar = iw / ih
    if bw / bh > ar:          # box wider than image -> height-limited
        h = bh
        w = bh * ar
    else:                     # width-limited
        w = bw
        h = bw / ar
    left = bl + (bw - w) / 2.0
    top = bt + (bh - h) / 2.0
    return left, top, w, h


def add_picture(slide, img_path, bl, bt, bw, bh):
    l, t, w, h = fit(img_path, bl, bt, bw, bh)
    slide.shapes.add_picture(img_path, Inches(l), Inches(t),
                             Inches(w), Inches(h))


def add_text(slide, text, l, t, w, h, size, color, bold=False,
             align=PP_ALIGN.LEFT, font=BODY_FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font
    f.color.rgb = color
    return tb


def title(slide, text, size=30):
    add_text(slide, text, 0.55, 0.32, 12.2, 0.95, size, INK,
             bold=True, font=TITLE_FONT)


def oneliner(slide, text, color=MUTED, bold=False, size=15):
    """A single short line pinned near the bottom, centered."""
    add_text(slide, text, 0.7, 6.78, 11.93, 0.5, size, color,
             bold=bold, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_number(slide, n):
    add_text(slide, str(n), 12.5, 7.02, 0.6, 0.35, 12, MUTED,
             align=PP_ALIGN.RIGHT)


def new(n):
    s = prs.slides.add_slide(BLANK)
    if n > 1:
        slide_number(s, n)
    return s


# layout constants for content slides
TOP = 1.35          # figure area starts below title
BOT_FULL = 6.85     # figure area bottom when no caption
BOT_CAP = 6.65      # figure area bottom when there's a one-liner
LEFT = 0.55
RIGHT_W = 12.23     # usable width

# ====================================================================
# 1. TITLE
# ====================================================================
s = new(1)
# left: text block, right: real two-channel hero image
add_text(s, "Measuring how proteins sense membrane curvature",
         0.7, 2.05, 6.4, 2.2, 34, INK, bold=True, font=TITLE_FONT)
add_text(s, "A deep-learning detector trained on a physics-calibrated "
            "simulator of our confocal microscope.",
         0.72, 4.35, 6.2, 1.4, 18, MUTED, font=BODY_FONT)
add_text(s, "Summer research talk", 0.72, 5.6, 6.0, 0.5, 14, MUTED)
add_picture(s, os.path.join(FIG, "hero_two_channel.png"),
            7.35, 1.0, 5.45, 5.5)

# ====================================================================
# 2. CURVATURE SENSING AS alpha
# ====================================================================
s = new(2)
title(s, "Curvature sensing as a single number, α")
half = (RIGHT_W - 0.4) / 2.0
add_picture(s, os.path.join(ASSET, "curvature_sensing.png"),
            LEFT, TOP, half, BOT_CAP - TOP)
add_picture(s, os.path.join(FIG, "sorting_curve_endophilin.png"),
            LEFT + half + 0.4, TOP, half, BOT_CAP - TOP)
oneliner(s, "Protein density ∝ d^(α−2)   •   "
            "α = 2: no sensing   •   α < 2: senses curvature",
         size=16)

# ====================================================================
# 3. TRAIN ON A CALIBRATED SIMULATOR
# ====================================================================
s = new(3)
title(s, "Train on a calibrated simulator, not real data")
add_picture(s, os.path.join(FIG, "real_vs_sim_tiles.png"),
            LEFT, TOP, RIGHT_W, BOT_CAP - TOP)
oneliner(s, "Real data has no ground truth; the calibrated simulator gives "
            "perfectly labeled training images.")

# ====================================================================
# 4. CALIBRATION: MATCH IMAGE STATISTICS  (title only)
# ====================================================================
s = new(4)
title(s, "Calibration: match image statistics, detection-free")
add_picture(s, os.path.join(FIG, "calibration_stats.png"),
            LEFT, TOP + 0.4, RIGHT_W, BOT_FULL - TOP - 0.8)

# ====================================================================
# 5. BOOTSTRAP RE-FITS
# ====================================================================
s = new(5)
title(s, "Validated across 100 bootstrap re-fits")
add_picture(s, os.path.join(FIG, "bootstrap_distributions.png"),
            LEFT, TOP, RIGHT_W, BOT_CAP - TOP)
oneliner(s, "Gain and excess-noise are degenerate — only their product "
            "is constrained, not each alone.")

# ====================================================================
# 6. U-NET vs CMEAnalysis  (promote r2 headline)
# ====================================================================
s = new(6)
title(s, "U-Net vs CMEAnalysis on real endophilin")
add_text(s, "r²  0.08 → 0.58  (7× better fit)   •   "
            "~3× more usable spots",
         0.55, 1.18, 12.2, 0.55, 19, BLUE, bold=True,
         align=PP_ALIGN.CENTER)
add_picture(s, os.path.join(FIG, "sorting_curve_compare.png"),
            LEFT, 1.9, RIGHT_W, BOT_FULL - 1.9)

# ====================================================================
# 7. EGFP NEGATIVE CONTROL
# ====================================================================
s = new(7)
title(s, "EGFP negative control reads as false sensing")
add_picture(s, os.path.join(FIG, "egfp_bias.png"),
            LEFT, TOP, RIGHT_W, BOT_CAP - TOP)
oneliner(s, "EGFP does not bind membranes — its slope should be 0. "
            "The detector reports sensing anyway.")

# ====================================================================
# 8. CANDIDATE CAUSES (not a verdict)
# ====================================================================
s = new(8)
title(s, "Candidate causes we're testing")
half = (RIGHT_W - 0.4) / 2.0
add_picture(s, os.path.join(FIG, "training_prior.png"),
            LEFT, TOP, half, BOT_CAP - TOP)
add_picture(s, os.path.join(FIG, "diameter_stratified.png"),
            LEFT + half + 0.4, TOP, half, BOT_CAP - TOP)
oneliner(s, "Possibilities under investigation: training distribution  •  "
            "resolution loss on small / dim spots  •  acquisition "
            "(lipid PMT lowered per sample, 750→580 V)", size=14)

# ====================================================================
# 9. PLAN 1: balanced data + full-resolution models
# ====================================================================
s = new(9)
title(s, "Plan 1: balanced data + full-resolution models")
add_picture(s, os.path.join(ASSET, "architectures.png"),
            LEFT, TOP, RIGHT_W, BOT_CAP - TOP)
oneliner(s, "Randomize curvature per spot (no learnable prior); test models "
            "that keep full resolution instead of shrinking the image.")

# ====================================================================
# 10. PLAN 2: condition on the DLS size prior
# ====================================================================
s = new(10)
title(s, "Plan 2: condition the detector on the DLS size prior")
# wide schematic on top, narrow null-result plot tucked beside it
left_w = RIGHT_W * 0.62
right_w = RIGHT_W - left_w - 0.4
add_picture(s, os.path.join(ASSET, "dls_conditioning.png"),
            LEFT, TOP, left_w, BOT_CAP - TOP)
add_picture(s, os.path.join(FIG, "film_null.png"),
            LEFT + left_w + 0.4, TOP, right_w, BOT_CAP - TOP)
oneliner(s, "FiLM on the plain size distribution was a null result; next is "
            "the d⁶-weighted input.")

# ====================================================================
# 11. PIPELINE + BENCHMARK + GOAL
# ====================================================================
s = new(11)
title(s, "Pipeline, benchmark, and goal")
add_picture(s, os.path.join(ASSET, "pipeline.png"),
            LEFT, TOP, RIGHT_W, BOT_CAP - TOP)
oneliner(s, "Out-of-the-box Spotiflow under-detects our small liposomes "
            "— an instrument mismatch.")

out = os.path.join(HERE, "deck.pptx")
prs.save(out)
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides")
